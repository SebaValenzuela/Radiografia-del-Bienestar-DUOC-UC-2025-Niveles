from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

def rellenar_template(template_path: str, output_path: str, placeholders: dict):
    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text in placeholders:
                    # Si el placeholder tiene una imagen en BytesIO
                    image_stream = placeholders[text]
                    # Limpiar el texto del placeholder
                    shape.text_frame.clear()
                    # Insertar la imagen en la misma posición y tamaño
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

    prs.save(output_path)
